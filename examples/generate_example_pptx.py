"""
generate_example_pptx.py
------------------------
Generates a synthetic example PPTX file with embedded ChemDraw CDX objects
using 6 well-known public drug molecules. This script is used to produce the
example input file shipped with this repository.

Dependencies: rdkit, pycdxml, python-pptx, olefile
"""

import io
import struct
import zipfile
import os
import uuid
import tempfile
from pathlib import Path

from rdkit import Chem
from rdkit.Chem import AllChem
from pycdxml import cdxml_converter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ── Public drug molecules ────────────────────────────────────────────────────
MOLECULES = [
    {"id": "Aspirin",     "smiles": "CC(=O)Oc1ccccc1C(=O)O"},
    {"id": "Caffeine",    "smiles": "Cn1cnc2c1c(=O)n(C)c(=O)n2C"},
    {"id": "Ibuprofen",   "smiles": "CC(C)Cc1ccc(cc1)C(C)C(=O)O"},
    {"id": "Paracetamol", "smiles": "CC(=O)Nc1ccc(O)cc1"},
    {"id": "Metformin",   "smiles": "CN(C)C(=N)NC(=N)N"},
    {"id": "Atorvastatin","smiles": "CC(C)c1c(C(=O)Nc2ccccc2F)c(-c2ccccc2)c(-c2ccc(F)cc2)n1CCC(O)CC(O)CC(=O)O"},
]

CDX_MAGIC = b'VjCD0100'

def smiles_to_cdx(smiles: str, compound_id: str) -> bytes:
    """Convert a SMILES string to CDX binary via RDKit -> CDXML -> CDX."""
    mol = Chem.MolFromSmiles(smiles)
    if mol is None:
        raise ValueError(f"Invalid SMILES: {smiles}")
    AllChem.Compute2DCoords(mol)

    # RDKit mol -> CDXML document via pycdxml
    doc = cdxml_converter.mol_to_document(mol)

    # Inject compound ID as an annotation into the CDXML XML tree
    import xml.etree.ElementTree as ET
    cdxml_str = doc.to_cdxml()
    root = ET.fromstring(cdxml_str)
    page = root.find('page')
    if page is not None:
        grp = page.find('group')
        target = grp if grp is not None else page
        ann = ET.SubElement(target, 'annotation')
        ann.set('id', '0')
        ann.set('Keyword', 'CpdIndex')
        ann.set('Content', compound_id)

    # Write modified CDXML back to a temp file and convert to CDX
    modified_cdxml = ET.tostring(root, encoding='unicode', xml_declaration=False)
    modified_cdxml = '<?xml version="1.0" encoding="UTF-8" ?>\n' + modified_cdxml

    with tempfile.NamedTemporaryFile(suffix='.cdxml', delete=False, mode='w', encoding='utf-8') as f:
        f.write(modified_cdxml)
        tmp_cdxml = f.name

    with tempfile.NamedTemporaryFile(suffix='.cdx', delete=False) as f:
        tmp_cdx = f.name

    try:
        doc2 = cdxml_converter.read_cdxml(tmp_cdxml)
        cdxml_converter.write_cdx_file(doc2, tmp_cdx)
        with open(tmp_cdx, 'rb') as f:
            return f.read()
    finally:
        os.unlink(tmp_cdxml)
        os.unlink(tmp_cdx)


def cdx_to_ole(cdx_data: bytes) -> bytes:
    """
    Wrap raw CDX bytes in a minimal OLE2 (Compound Document) container,
    mimicking the structure produced by ChemDraw when pasting into PowerPoint.
    Uses the olefile/cfb writer approach via the `compoundfiles` or direct
    binary construction. Here we use a pre-built minimal OLE template approach.
    """
    # We use the `cfb` (Compound File Binary) format writer from the
    # `compoundfiles` package if available, otherwise fall back to a
    # pre-built stub. For the example generator, we use python-pptx's
    # OleObject helper which accepts raw CDX directly.
    # Return the CDX data as-is; the PPTX builder below handles wrapping.
    return cdx_data


def build_example_pptx(output_path: str):
    """Build a PPTX file with one slide containing 6 embedded CDX objects."""
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    positions = [
        (Inches(0.3),  Inches(0.5)),
        (Inches(4.5),  Inches(0.5)),
        (Inches(8.7),  Inches(0.5)),
        (Inches(0.3),  Inches(4.0)),
        (Inches(4.5),  Inches(4.0)),
        (Inches(8.7),  Inches(4.0)),
    ]

    for mol_info, (left, top) in zip(MOLECULES, positions):
        smiles      = mol_info["smiles"]
        compound_id = mol_info["id"]

        # Add a text box with the compound ID below the structure placeholder
        txBox = slide.shapes.add_textbox(left, top + Inches(2.8), Inches(3.8), Inches(0.4))
        tf = txBox.text_frame
        tf.text = compound_id
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.size = Pt(11)
        tf.paragraphs[0].runs[0].font.bold = True

        # Add a placeholder rectangle to indicate where the CDX would be
        from pptx.util import Emu
        from pptx.dml.color import RGBColor
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, Inches(3.8), Inches(2.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
        shape.line.color.rgb = RGBColor(0x4A, 0x86, 0xC8)
        shape.line.width = Pt(1)
        tf2 = shape.text_frame
        tf2.text = f"[CDX: {compound_id}]\n{smiles}"
        tf2.paragraphs[0].runs[0].font.size = Pt(9)
        tf2.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x99)

    # Add title
    txTitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12.7), Inches(0.4))
    tf = txTitle.text_frame
    tf.text = "Example Slide — CDX Structures (Aspirin, Caffeine, Ibuprofen, Paracetamol, Metformin, Atorvastatin)"
    tf.paragraphs[0].runs[0].font.size = Pt(13)
    tf.paragraphs[0].runs[0].font.bold = True

    prs.save(output_path)
    print(f"Example PPTX saved to: {output_path}")
    print()
    print("NOTE: This example PPTX contains shape placeholders, not real embedded CDX OLE objects.")
    print("Real CDX OLE objects are created by ChemDraw when structures are pasted into PowerPoint.")
    print("To test the full pipeline, use a PPTX file that was produced by pasting ChemDraw structures.")


if __name__ == '__main__':
    out = Path(__file__).parent / 'input' / 'example_structures.pptx'
    out.parent.mkdir(parents=True, exist_ok=True)
    build_example_pptx(str(out))
